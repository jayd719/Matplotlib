import pandas as pd
from pptx import Presentation
from pptx.util import Inches

# Create a presentation
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Virtue Ethics in Technology: Strengths, Challenges, and Real-World Cases"
subtitle.text = (
    "An Analysis of Virtue Ethics in Computer Science and Ethical Decision-Making"
)

# Slide 1: Introduction
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = (
    "Virtue ethics emphasizes moral character rather than strict rules or consequences.\n"
    "This approach can lead to ethical ambiguity in technology fields like AI, cybersecurity, and data privacy.\n"
    "While virtues such as honesty, fairness, and responsibility are essential, they alone may not provide clear guidance."
)

# Slide 2: Challenges of Relying Solely on Virtue Ethics
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Challenges of Relying Solely on Virtue Ethics"
content.text = (
    "1. Lack of Clear Ethical Standards - Ethics based on individual virtues leads to inconsistent decision-making.\n"
    "2. Subjectivity & Moral Relativism - Different interpretations of virtues can lead to conflicting ethical choices.\n"
    "3. Decision-Making Under Pressure - No structured framework makes quick ethical decisions difficult.\n"
    "4. Ethical Blind Spots & Rationalization - Can be used to justify unethical actions.\n"
    "5. Lack of Accountability - No enforcement mechanisms for ethical behavior in companies."
)

# Slide 3: Real-World Example: Facebook & Cambridge Analytica
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Case Study 1: Facebook & Cambridge Analytica"
content.text = (
    "- Facebook allowed massive user data collection without informed consent.\n"
    "- Virtue of innovation prioritized over privacy and responsibility.\n"
    "- Lack of clear ethical regulations led to data misuse.\n"
    "- GDPR and other regulations emerged to address such ethical gaps."
)

# Slide 4: Real-World Example: Google & Project Maven
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Case Study 2: Google & Project Maven"
content.text = (
    "- Google worked with the Pentagon on AI for drone footage analysis.\n"
    "- Employees protested, arguing the project conflicted with ethical virtues.\n"
    "- Internal ethical debates led to Google discontinuing the project.\n"
    "- Highlights the need for structured ethical guidelines in AI ethics."
)

# Slide 5: The Need for a Hybrid Approach
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "The Need for a Hybrid Approach"
content.text = (
    "- Combining Virtue Ethics with other frameworks provides a balanced ethical foundation.\n"
    "- Deontology: Clear rules (e.g., 'Do not steal data').\n"
    "- Utilitarianism: Focus on outcomes and consequences.\n"
    "- Contractualism: Respect for individual rights and fairness.\n"
    "- Combining these approaches ensures both ethical integrity and accountability."
)

# Slide 6: Conclusion
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = (
    "- Virtue ethics promotes ethical awareness, but it is not enough on its own.\n"
    "- Structured ethical policies and regulations are necessary.\n"
    "- A hybrid ethical approach ensures responsible technology development.\n"
    "- Ethics in tech must balance innovation, societal well-being, and accountability."
)

# Save the presentation
pptx_filename = "Virtue_Ethics_Technology.pptx"
prs.save(pptx_filename)

# Provide the file to the user
pptx_filename
